# -*- coding: utf-8 -*-
"""
Genera Presentacion_TFG.pptx clonando la plantilla del Departamento de Historia
del Arte (UM) -- la misma que la presentacion 'Memorias artificiales' (10).

Paleta exacta: fondo #FFFEF5, rojo #E73C23, azul marino #03293B.
Logos oficiales extraidos del PDF de ejemplo en presentacion/assets/.

Diseño pensado para transmitir a un publico EXPERTO y a uno LEGO:
cada diapositiva tiene una idea-titular accesible + terminos clave en negrita.
Las notas del orador (panel de notas) llevan la doble lectura.

Ejecutar desde la carpeta del repo:  python3 presentacion/generar_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(REPO, "presentacion", "assets")
OUT = os.path.join(REPO, "presentacion", "Presentacion_TFG.pptx")

CREMA = RGBColor(0xFF, 0xFE, 0xF5)
ROJO = RGBColor(0xE7, 0x3C, 0x23)
AZUL = RGBColor(0x03, 0x29, 0x3B)
GRIS = RGBColor(0x44, 0x44, 0x44)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

TITULO = "LA VISIBILIDAD DEL CINE ARGENTINO EN ESPAÑA (1980-2000)"
FONT = "Calibri"

LOGO_UMU = os.path.join(ASSETS, "logo_umu.png")
LOGO_FAC = os.path.join(ASSETS, "logo_facultad.png")
LOGO_COR = os.path.join(ASSETS, "logo_corazon.png")

IMG = {
    "hogar": "Tranquilidad y felicidad del hogar de Alicia frente a la verdad del exterior.png",
    "madres": "Alicia contemplando la marcha de las Madres de Mayo .png",
    "ana": "Ana le cuenta lo que la violencia que sufrió a Alicia.png",
    "paisaje": "Plano del paisaje en San Luis.png",
    "tren": "Contraste entre el mundo del avance (tren) y mundo rural.png",
    "mario": "Mario imparte clase a los niños locales.png",
    "pacifico": "Cartel de Nueva York y Pacifico en el video de Hache.png",
    "aeropuerto": "Dante recibe a Martin y a Hache en el aeropuerto.png",
    "sobredosis": "Hache sufriendo los efectos de la sobredosis.png",
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = CREMA
    bg.line.fill.background(); bg.shadow.inherit = False
    bg._element.addprevious(bg._element)  # keep at back (no-op safe)
    return s


def _set_runs(p, segments, size, color, italic=False):
    """segments: str o lista de (texto, bold)."""
    if isinstance(segments, str):
        segments = [(segments, False)]
    for (t, bold) in segments:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = FONT


def text(s, x, y, w, h, paragraphs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         size=18, color=AZUL, sp_after=6, line=None, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        if line:
            p.line_spacing = line
        seg = para.get("seg") if isinstance(para, dict) else para
        psize = para.get("size", size) if isinstance(para, dict) else size
        pcolor = para.get("color", color) if isinstance(para, dict) else color
        pital = para.get("italic", italic) if isinstance(para, dict) else italic
        _set_runs(p, seg, psize, pcolor, pital)
    return tb


def title(s, txt, color=AZUL, size=28):
    text(s, Inches(0.55), Inches(0.35), Inches(12.2), Inches(1.1),
         [txt], size=size, color=color, line=1.0)


def footer(s, kind="content"):
    # linea horizontal
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(6.78), SW, Pt(1.6))
    ln.fill.solid(); ln.fill.fore_color.rgb = AZUL; ln.line.fill.background(); ln.shadow.inherit = False
    # logo UMU abajo-izquierda
    s.shapes.add_picture(LOGO_UMU, Inches(0.45), Inches(6.95), height=Inches(0.42))
    if kind == "content":
        text(s, Inches(3.4), Inches(6.92), Inches(9.5), Inches(0.5),
             [TITULO], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             size=11.5, color=ROJO, sp_after=0)


def callout(s, x, y, w, h, segments):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = CREMA
    box.line.color.rgb = ROJO; box.line.width = Pt(1.75); box.shadow.inherit = False
    try:
        box.adjustments[0] = 0.18
    except Exception:
        pass
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    _set_runs(p, segments, 16, AZUL)
    return box


def arrow(s, x, y, w=Inches(0.9), h=Inches(0.32)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = ROJO; a.line.fill.background(); a.shadow.inherit = False
    return a


def pic_border(s, key, x, y, boxw, boxh, color=AZUL):
    """Inserta la imagen escalada para caber en boxw x boxh, centrada, con borde."""
    path = os.path.join(REPO, IMG[key])
    if not os.path.exists(path):
        print("  ! Falta:", IMG[key]); return None
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(boxw / iw, boxh / ih)
    w = int(iw * scale); h = int(ih * scale)
    px = x + (boxw - w) // 2
    py = y + (boxh - h) // 2
    pic = s.shapes.add_picture(path, px, py, width=w, height=h)
    pic.line.color.rgb = color; pic.line.width = Pt(3)
    return pic


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


# ============================================================ 1. PORTADA
s = slide()
# corazon UMU decorativo, arriba-derecha sangrando
s.shapes.add_picture(LOGO_COR, Inches(9.7), Inches(-0.4), width=Inches(4.6))
text(s, Inches(0.7), Inches(1.35), Inches(8.7), Inches(2.6), [
    {"seg": [("LA VISIBILIDAD DEL CINE ARGENTINO ", True)], "size": 33, "color": ROJO},
    {"seg": [("EN ESPAÑA (1980-2000)", True)], "size": 33, "color": ROJO},
], line=1.05, sp_after=2)
text(s, Inches(0.72), Inches(3.5), Inches(8.6), Inches(0.8),
     [{"seg": [("Memoria, exilio y coproducciones transatlánticas", False)], "italic": True}],
     size=20, color=AZUL)
text(s, Inches(0.72), Inches(4.9), Inches(8.6), Inches(1.7), [
    {"seg": [("Facundo Nicolás Vallejos Ribles", True)], "size": 19, "color": AZUL},
    {"seg": [("TFG · Grado en Historia del Arte", False)], "size": 15, "color": AZUL},
    {"seg": [("Prof. D. Joaquín Cánovas Belchí", False)], "size": 15, "color": AZUL},
    {"seg": [("Curso 2025/2026 · Convocatoria junio-julio", False)], "size": 15, "color": AZUL},
], sp_after=6)
# pie portada: linea + logos + departamento
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(6.78), SW, Pt(1.6))
ln.fill.solid(); ln.fill.fore_color.rgb = AZUL; ln.line.fill.background(); ln.shadow.inherit = False
s.shapes.add_picture(LOGO_UMU, Inches(0.45), Inches(6.95), height=Inches(0.42))
s.shapes.add_picture(LOGO_FAC, Inches(3.55), Inches(6.86), height=Inches(0.58))
text(s, Inches(8.3), Inches(6.92), Inches(4.6), Inches(0.5),
     ["Departamento de Historia del Arte"], align=PP_ALIGN.RIGHT,
     anchor=MSO_ANCHOR.MIDDLE, size=12, color=AZUL, sp_after=0)
notes(s, "Saludo breve y tranquilo. 'Buenos dias. Muchas gracias, miembros del tribunal. "
         "Voy a defender mi TFG sobre la presencia del cine argentino en España entre 1980 y 2000.' "
         "No leas el titulo: ya esta en pantalla. Mira al tribunal.")

# ============================================================ 2. ÍNDICE
s = slide()
text(s, Inches(9.8), Inches(0.4), Inches(3.0), Inches(0.8),
     [{"seg": [("ÍNDICE", True)]}], align=PP_ALIGN.RIGHT, size=30, color=AZUL)
idx = [
    ([("INTRODUCCIÓN", True)], ROJO, 0.0),
    ([("Justificación y estado de la cuestión", False)], AZUL, 0.6),
    ([("Hipótesis y objetivos", False)], AZUL, 0.6),
    ([("Metodología", False)], AZUL, 0.6),
    ([("CONTENIDO", True)], ROJO, 0.0),
    ([("Contexto: memoria, exilio y coproducciones", False)], AZUL, 0.6),
    ([("Tres casos de estudio: ", False), ("La historia oficial, Un lugar en el mundo, Martín (Hache)", True)], AZUL, 0.6),
    ([("Recepción en España", False)], AZUL, 0.6),
    ([("CONCLUSIONES", True)], ROJO, 0.0),
]
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
for i, (seg, col, ind) in enumerate(idx):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10); p.level = 0
    p.alignment = PP_ALIGN.LEFT
    if ind:
        seg = [("        ", False)] + seg
    _set_runs(p, seg, 19 if col == ROJO else 17, col)
footer(s)
notes(s, "Hoja de ruta en 15 segundos. 'La presentacion sigue la estructura clasica: "
         "introduccion, contenido con tres casos de cine, y conclusiones.'")

# ============================================================ 3. ESTADO DE LA CUESTIÓN
s = slide()
title(s, [("JUSTIFICACIÓN Y ESTADO DE LA CUESTIÓN", True)], color=ROJO, size=26)
ejes = [
    ("MEMORIA", "del trauma de la dictadura", Inches(1.2), Inches(2.2)),
    ("COPRODUCCIONES", "hispano-argentinas", Inches(7.0), Inches(2.2)),
    ("EXILIO", "como mediación cultural", Inches(1.2), Inches(4.05)),
    ("RECEPCIÓN", "crítica e institucional", Inches(7.0), Inches(4.05)),
]
for name, sub, x, yy in ejes:
    arrow(s, x, yy + Inches(0.18), w=Inches(0.85), h=Inches(0.3))
    text(s, x + Inches(1.0), yy - Inches(0.05), Inches(4.5), Inches(1.2),
         [{"seg": [(name, True)], "size": 21, "color": AZUL},
          {"seg": [(sub, False)], "size": 13, "color": GRIS}],
         sp_after=2, line=1.0)
# franja inferior con el hueco
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.75), Inches(11.3), Inches(0.85))
box.fill.solid(); box.fill.fore_color.rgb = CREMA; box.line.color.rgb = ROJO; box.line.width = Pt(1.75); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_set_runs(p, [("La bibliografía aborda el tema desde lo político y social. El hueco que llena este TFG: el ", False),
              ("análisis FORMAL", True), (" del cine como objeto visual.", False)], 15, AZUL)
footer(s)
notes(s, "EXPERTO: situa el debate (Amado y la 'imagen justa', Jensen sobre el exilio, Elena con los datos de circulacion). "
         "LEGO: 'Mucho se ha escrito sobre la politica y la historia de estas peliculas, pero casi nadie las ha mirado "
         "como lo que son: imagenes. Eso es lo que aporto, mirarlas con ojos de historiador del arte.'")

# ============================================================ 4. HIPÓTESIS Y OBJETIVOS
s = slide()
title(s, [("HIPÓTESIS Y OBJETIVOS", True)], color=ROJO, size=26)
hyp = [
    ("1.", "MEMORIA", "El cine argentino tuvo una recepción singular en España por tratar la represión política, reconocible desde la propia dictadura y transición españolas."),
    ("2.", "EXILIO", "El exiliado argentino no fue un sujeto pasivo, sino un mediador cultural que tejió redes de intercambio artístico entre los dos países."),
    ("3.", "COPRODUCCIONES", "Fueron una de las vías más eficaces de visibilidad del cine argentino, con sus posibilidades pero también con sus tensiones."),
]
y = Inches(1.7)
for num, key, body in hyp:
    text(s, Inches(0.8), y, Inches(0.7), Inches(1.0), [{"seg": [(num, True)], "size": 26, "color": ROJO}])
    text(s, Inches(1.6), y, Inches(11.0), Inches(1.3), [
        {"seg": [(key, True)], "size": 18, "color": AZUL},
        {"seg": [(body, False)], "size": 15.5, "color": GRIS},
    ], sp_after=2, line=1.02)
    y += Inches(1.5)
footer(s)
notes(s, "Estas tres hipotesis son la columna vertebral; en conclusiones se 'resuelven'. Apréndelas de memoria. "
         "LEGO: traduce cada una a una frase ('¿por que gusto aqui? por la memoria; ¿quien tendio el puente? los exiliados; "
         "¿como llego? por las coproducciones').")

# ============================================================ 5. METODOLOGÍA
s = slide()
title(s, [("METODOLOGÍA", True)], color=ROJO, size=26)
text(s, Inches(0.9), Inches(2.9), Inches(2.2), Inches(1.2),
     [{"seg": [("MÉTODO", True)]}], size=26, color=ROJO, anchor=MSO_ANCHOR.MIDDLE)
cy = [Inches(1.55), Inches(3.0), Inches(4.45)]
metodo = [
    [("Análisis FORMAL: ", True), ("las películas como textos visuales (puesta en escena, espacio, luz, encuadre, montaje).", False)],
    [("Aproximación HISTÓRICO-CULTURAL: ", True), ("contexto de dictadura, exilio y relaciones España-Argentina.", False)],
    [("Fuentes: ", True), ("Amado, Sarlo, Jensen, Elena y González; pautas de Cánovas Belchí y Aliaga Cárceles. Citas en Chicago-Deusto.", False)],
]
for yy, seg in zip(cy, metodo):
    arrow(s, Inches(3.2), yy + Inches(0.35), w=Inches(0.75), h=Inches(0.28))
    callout(s, Inches(4.1), yy, Inches(8.3), Inches(1.05), seg)
footer(s)
notes(s, "EXPERTO: subraya que combinas analisis formal e historico-cultural y que lees el film como texto visual. "
         "LEGO: 'Veo las peliculas con lupa -como se coloca la camara, que se ilumina, que espacio se muestra- y "
         "ademas las situo en su momento historico.'")

# ============================================================ 6. CONTEXTO
s = slide()
title(s, [("CONTEXTO: DOS SOCIEDADES, UNA MEMORIA COMPARTIDA", True)], color=AZUL, size=23)
# Argentina
text(s, Inches(0.8), Inches(1.75), Inches(5.6), Inches(0.5), [{"seg": [("ARGENTINA", True)]}], size=18, color=ROJO)
text(s, Inches(0.8), Inches(2.35), Inches(5.6), Inches(4.0), [
    {"seg": [("Dictadura (1976-1983): ", True), ("secuestro, tortura y desaparición.", False)], "size": 15},
    {"seg": [("El exilio se vuelve necesidad de supervivencia ", False), ("(Briski, Guevara, Alterio).", True)], "size": 15},
], color=GRIS, sp_after=10, line=1.05)
# España
text(s, Inches(7.0), Inches(1.75), Inches(5.5), Inches(0.5), [{"seg": [("ESPAÑA", True)]}], size=18, color=ROJO)
text(s, Inches(7.0), Inches(2.35), Inches(5.5), Inches(4.0), [
    {"seg": [("Muerte de Franco (1975) ", True), ("y transición democrática.", False)], "size": 15},
    {"seg": [("Madrid y Barcelona, ", False), ("espacios receptivos para el exilio argentino.", True)], "size": 15},
], color=GRIS, sp_after=10, line=1.05)
# franja union
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(5.4), Inches(10.7), Inches(1.0))
box.fill.solid(); box.fill.fore_color.rgb = AZUL; box.line.fill.background(); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_set_runs(p, [("Dos países que salían de la violencia de Estado se reconocen en las mismas imágenes.", True)], 16, CREMA)
footer(s)
notes(s, "Idea clave y muy accesible: el cine argentino no llego a un pais cualquiera, sino a otro que tambien "
         "venia de una dictadura. Por eso 'resonaba'. No narres toda la historia politica: una frase por lado.")

# ============================================================ 7. MEMORIA 80 vs 90
s = slide()
title(s, [("LA MEMORIA EN PANTALLA: DE LOS 80 A LOS 90", True)], color=AZUL, size=23)
text(s, Inches(0.8), Inches(1.85), Inches(5.6), Inches(0.5), [{"seg": [("AÑOS 80", True)]}], size=18, color=ROJO)
text(s, Inches(0.8), Inches(2.45), Inches(5.6), Inches(4.0), [
    {"seg": [("Realismo afectivo", True), (" y «imaginación melodramática».", False)], "size": 15},
    {"seg": [("Vocación didáctica, ", False), ("relatos claros y directos.", False)], "size": 15},
    {"seg": [("El espacio familiar ", True), ("como espejo de la sociedad.", False)], "size": 15},
], color=GRIS, sp_after=9, line=1.05)
text(s, Inches(7.0), Inches(1.85), Inches(5.5), Inches(0.5), [{"seg": [("AÑOS 90 · NUEVO CINE ARGENTINO", True)]}], size=15, color=ROJO)
text(s, Inches(7.0), Inches(2.45), Inches(5.5), Inches(4.0), [
    {"seg": [("Ambigüedad, fragmento ", True), ("y distancia crítica.", False)], "size": 15},
    {"seg": [("Richard: ", True), ("la memoria es un territorio en disputa.", False)], "size": 15},
    {"seg": [("Sarlo: ", True), ("el riesgo del «giro subjetivo».", False)], "size": 15},
], color=GRIS, sp_after=9, line=1.05)
footer(s)
notes(s, "EXPERTO: 'imagen justa' (Amado), giro subjetivo (Sarlo), memoria como disputa (Richard). "
         "LEGO: 'En los 80 el cine explicaba el horror de forma clara y emotiva, casi pedagogica. En los 90 se vuelve "
         "mas sugerente, deja preguntas abiertas.' Esto prepara el contraste entre los tres casos.")

# ============================================================ 8. CASO 1
s = slide()
title(s, [("CASO 1 · ", True), ("LA HISTORIA OFICIAL", True), (" (Puenzo, 1985)", False)], color=AZUL, size=23)
text(s, Inches(0.8), Inches(1.5), Inches(11.6), Inches(0.6),
     [{"seg": [("Idea clave: el espacio como ", False), ("CEGUERA", True), (" — el hogar burgués, cerrado a la verdad de fuera.", False)]}],
     size=15, color=ROJO, italic=False)
text(s, Inches(3.0), Inches(2.15), Inches(0.6), Inches(0.4), [{"seg": [("INTERIOR", True)]}], size=12, color=ROJO)
text(s, Inches(9.0), Inches(2.15), Inches(2.0), Inches(0.4), [{"seg": [("EXTERIOR", True)]}], size=12, color=ROJO)
pic_border(s, "hogar", Inches(0.9), Inches(2.55), Inches(5.6), Inches(3.2), color=AZUL)
pic_border(s, "madres", Inches(6.85), Inches(2.55), Inches(5.6), Inches(3.2), color=ROJO)
text(s, Inches(0.9), Inches(5.95), Inches(5.6), Inches(0.7),
     [{"seg": [("El hogar acomodado de Alicia: aparente felicidad.", False)]}], size=12.5, color=AZUL, align=PP_ALIGN.CENTER)
text(s, Inches(6.85), Inches(5.95), Inches(5.6), Inches(0.7),
     [{"seg": [("Las Madres de Plaza de Mayo: la verdad del exterior.", False)]}], size=12.5, color=AZUL, align=PP_ALIGN.CENTER)
footer(s)
notes(s, "LEGO (argumento en 2 frases): una profesora de clase acomodada empieza a sospechar que su hija adoptiva "
         "puede ser hija de desaparecidos. EXPERTO: Puenzo viene de la publicidad; puesta en escena muy controlada; "
         "interiores calidos vs exteriores frios = ceguera social; la transformacion de Alicia como alegoria nacional. "
         "Dato: Oscar 1986. Matiz critico: el melodrama puede rozar la 'teoria de los dos demonios'.")

# ============================================================ 9. EXILIO
s = slide()
title(s, [("EL EXILIO: UN PUENTE ENTRE PAÍSES", True)], color=ROJO, size=25)
text(s, Inches(0.9), Inches(1.6), Inches(2.4), Inches(1.0), [{"seg": [("EXILIO", True)]}], size=24, color=ROJO, anchor=MSO_ANCHOR.MIDDLE)
cy = [Inches(1.55), Inches(3.0), Inches(4.45)]
seg3 = [
    [("No es migración económica: ", True), ("es huida de una violencia que amenaza la vida.", False)],
    [("«Figuras del destierro» (Rodríguez Marino): ", True), ("la pérdida y el retorno imposible se hacen lenguaje visual.", False)],
    [("El exiliado tiende puentes: ", True), ("Reflexiones de un salvaje (Vallejo, 1978) une franquismo y represión argentina.", False)],
]
for yy, seg in zip(cy, seg3):
    arrow(s, Inches(3.3), yy + Inches(0.35), w=Inches(0.75), h=Inches(0.28))
    callout(s, Inches(4.2), yy, Inches(8.2), Inches(1.05), seg)
footer(s)
notes(s, "LEGO: distingue exiliado (huye para salvar la vida) de emigrante economico. "
         "EXPERTO: 'figuras del destierro' de Rodriguez Marino; Vallejo rueda en España y equipara la represion "
         "franquista con la argentina (Don Quijote/Martin Fierro, los republicanos asesinados, montaje a lo Eisenstein). "
         "Si vas justo de tiempo, esta diapositiva es la mas recortable.")

# ============================================================ 10. COPRODUCCIONES
s = slide()
title(s, [("LAS COPRODUCCIONES: LA VÍA INDUSTRIAL", True)], color=ROJO, size=24)
text(s, Inches(0.9), Inches(1.55), Inches(2.6), Inches(1.0), [{"seg": [("COPRO-\nDUCCIÓN", True)]}], size=18, color=ROJO, anchor=MSO_ANCHOR.MIDDLE, line=1.0)
cy = [Inches(1.5), Inches(2.95), Inches(4.4)]
seg3 = [
    [("Más que una fórmula legal: ", True), ("estrategia de supervivencia frente al dominio de Hollywood.", False)],
    [("Convenio Hispano-Argentino de 1969: ", True), ("las películas son «nacionales» en ambos países (participación 30 %-70 %).", False)],
    [("Distintos modelos: ", True), ("Aristarain (autoral) frente a Piñeyro (industrial: Plata quemada, Kamchatka).", False)],
]
for yy, seg in zip(cy, seg3):
    arrow(s, Inches(3.5), yy + Inches(0.35), w=Inches(0.7), h=Inches(0.28))
    callout(s, Inches(4.3), yy, Inches(8.1), Inches(1.05), seg)
text(s, Inches(1.0), Inches(5.95), Inches(11.4), Inches(0.7),
     [{"seg": [("Tensión clave: la asimetría económica condiciona ", False), ("qué imagen de Argentina circula", True), (" en España.", False)]}],
     size=15, color=AZUL, align=PP_ALIGN.CENTER, italic=True)
footer(s)
notes(s, "LEGO: una coproduccion es 'poner dinero entre varios paises para hacer una pelicula'; eso le abre cines y "
         "ayudas a las que sola no llegaria. EXPERTO: Convenio de 1969 (segundo mas antiguo de España), doble "
         "nacionalidad; modelos autoral vs industrial; tension: el socio fuerte puede condicionar el relato.")

# ============================================================ 11. CASO 2
s = slide()
title(s, [("CASO 2 · ", True), ("UN LUGAR EN EL MUNDO", True), (" (Aristarain, 1992)", False)], color=AZUL, size=22)
text(s, Inches(0.8), Inches(1.5), Inches(11.6), Inches(0.6),
     [{"seg": [("Idea clave: el espacio ", False), ("UNITARIO", True), (" — todavía hay un lugar que defender.", False)]}],
     size=15, color=ROJO)
pic_border(s, "paisaje", Inches(0.9), Inches(2.35), Inches(5.6), Inches(3.3), color=AZUL)
pic_border(s, "tren", Inches(6.85), Inches(2.35), Inches(5.6), Inches(3.3), color=ROJO)
text(s, Inches(0.9), Inches(5.75), Inches(5.6), Inches(0.9),
     [{"seg": [("El paisaje árido de San Luis, casi un personaje (diálogo con el western).", False)]}], size=12.5, color=AZUL, align=PP_ALIGN.CENTER, line=1.0)
text(s, Inches(6.85), Inches(5.75), Inches(5.6), Inches(0.9),
     [{"seg": [("El tren: el progreso que avanza pero arrasa lo que no se adapta.", False)]}], size=12.5, color=AZUL, align=PP_ALIGN.CENTER, line=1.0)
footer(s)
notes(s, "LEGO: un maestro rural resiste eticamente en un pueblo perdido; su hijo recuerda esa infancia. "
         "EXPERTO: estetica de western (paisaje = aislamiento y resistencia); el tren como progreso destructor; "
         "un tercio de la pelicula son trayectos que 'cosen' un mundo unitario = 'tener un lugar en el mundo'. "
         "Dato: Concha de Oro (San Sebastian), Goya y +500.000 espectadores en España.")

# ============================================================ 12. CASO 3
s = slide()
title(s, [("CASO 3 · ", True), ("MARTÍN (HACHE)", True), (" (Aristarain, 1997)", False)], color=AZUL, size=22)
text(s, Inches(0.8), Inches(1.5), Inches(11.6), Inches(0.6),
     [{"seg": [("Idea clave: el espacio ", False), ("FRAGMENTADO", True), (" = identidad fragmentada. La memoria, por omisión.", False)]}],
     size=15, color=ROJO)
pic_border(s, "aeropuerto", Inches(0.9), Inches(2.35), Inches(5.6), Inches(3.3), color=AZUL)
pic_border(s, "pacifico", Inches(6.85), Inches(2.35), Inches(5.6), Inches(3.3), color=ROJO)
text(s, Inches(0.9), Inches(5.75), Inches(5.6), Inches(0.9),
     [{"seg": [("Madrid y Buenos Aires casi indiferenciados: pocos exteriores.", False)]}], size=12.5, color=AZUL, align=PP_ALIGN.CENTER, line=1.0)
text(s, Inches(6.85), Inches(5.75), Inches(5.6), Inches(0.9),
     [{"seg": [("«Pacífico» sobre Nueva York: la nostalgia inventa horizontes imposibles.", False)]}], size=12.5, color=AZUL, align=PP_ALIGN.CENTER, line=1.0)
footer(s)
notes(s, "LEGO: un cineasta argentino exiliado en Madrid recibe a su hijo adolescente tras una sobredosis; "
         "ninguno termina de pertenecer a ningun sitio. EXPERTO: predominio del dialogo (cine de autor europeo); "
         "espacio fragmentado = identidad rota; se OMITE la dictadura (memoria latente en el silencio); "
         "el titulo: la 'hache' muda, el tiempo en suspenso; el brindis final 'por Martin' restituye el nombre.")

# ============================================================ 13. RECEPCIÓN
s = slide()
title(s, [("RECEPCIÓN EN ESPAÑA", True)], color=ROJO, size=26)
text(s, Inches(0.85), Inches(1.8), Inches(11.6), Inches(4.5), [
    {"seg": [("Festivales y crítica: ", True), ("se lee el cine argentino como una cinematografía con identidad propia. Elena llama a Aristarain el «buque insignia».", False)], "size": 16},
    {"seg": [("Premios Goya: ", True), ("Un lugar en el mundo, Cenizas del paraíso y el Goya de Cecilia Roth por Martín (Hache).", False)], "size": 16},
    {"seg": [("El matiz importante: ", True), ("éxito sobre todo CRÍTICO e INSTITUCIONAL. Salvo Un lugar en el mundo, la mayoría circuló en salas especializadas y versión original.", False)], "size": 16},
], color=GRIS, sp_after=14, line=1.05)
footer(s)
notes(s, "Honestidad academica (gusta al tribunal): no idealizar. Prestigio critico ≠ exito de taquilla. "
         "LEGO: 'triunfo en festivales y premios, pero no fue cine masivo, salvo una excepcion.'")

# ============================================================ 14. CONCLUSIONES
s = slide()
title(s, [("CONCLUSIONES", True)], color=ROJO, size=26)
conc = [
    [("El ESPACIO como herramienta de sentido: ", True), ("ceguera (La historia oficial), unidad (Un lugar en el mundo), fragmentación (Martín (Hache)).", False)],
    [("Dos modos de MEMORIA: ", True), ("realismo afectivo y didáctico en los 80 frente a ambigüedad y distanciamiento en los 90.", False)],
    [("El EXILIO como tejido conectivo ", True), ("que preparó el terreno de las futuras coproducciones.", False)],
    [("Hipótesis confirmadas: ", True), ("España ofreció una recepción singular por resonancia con su propio pasado.", False)],
]
y = Inches(1.7)
for i, seg in enumerate(conc, 1):
    text(s, Inches(0.8), y, Inches(0.7), Inches(0.9), [{"seg": [("%d." % i, True)], "size": 22, "color": ROJO}])
    text(s, Inches(1.55), y + Inches(0.02), Inches(11.0), Inches(1.1), [{"seg": seg, "size": 16, "color": GRIS}], line=1.02)
    y += Inches(1.18)
footer(s)
notes(s, "Cierra con fuerza y mirando al tribunal. La frase de oro: 'Entre 1980 y 2000 las imagenes cruzaron el "
         "Atlantico cargadas de memoria, pero tambien de nuevas posibilidades esteticas e industriales.' "
         "El hilo conductor de toda la charla ha sido EL ESPACIO.")

# ============================================================ 15. REFERENCIAS
s = slide()
title(s, [("REFERENCIAS BIBLIOGRÁFICAS Y DOCUMENTALES", True)], color=ROJO, size=22)
refs = [
    "Amado, Ana María. La imagen justa: cine argentino y política (1980-2007). Buenos Aires: Colihue, 2009.",
    "Campero, Agustín. Nuevo cine argentino: de Rapado a Historias extraordinarias. Buenos Aires: UNGS/BN, 2009.",
    "Elena, Alberto. «La difusión del cine latinoamericano en España: una aproximación cuantitativa». Madrid: AEHC, 1998.",
    "González, Leandro Ramiro. «Cruzando el Atlántico: cine argentino en España». Imagofagia 17 (2018): 41-70.",
    "Jensen, Silvina. Los exiliados. La lucha por los derechos humanos durante la dictadura. Buenos Aires: Sudamericana, 2010.",
    "Richard, Nelly. Fracturas de la memoria: arte y pensamiento crítico. Buenos Aires: Siglo XXI, 2007.",
    "Rodríguez Marino, Paula. Figuras del destierro. Narraciones del exilio en el cine argentino (1978-1988). Viedma: UNRN, 2013.",
    "Sarlo, Beatriz. Tiempo pasado. Cultura de la memoria y giro subjetivo. Buenos Aires: Siglo XXI, 2005.",
    "Thibaudeau, Pascale. «Martín (Hache): del exilio al regreso, el itinerario de una transmisión». ALHIM 6 (2003).",
    "Convenio Hispano-Argentino de Relaciones Cinematográficas. BOE n.º 238, 4 de octubre de 1969, 15550-15552.",
]
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.1)); tf = tb.text_frame; tf.word_wrap = True
for i, r in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6); p.line_spacing = 1.0
    _set_runs(p, [("–  " + r, False)], 12.5, GRIS)
footer(s)
notes(s, "No leas las referencias. Estan para mostrar solidez bibliografica. Si preguntan por una fuente concreta, "
         "aqui la tienes localizada. Estilo Chicago-Deusto, como exige el Grado.")

# ============================================================ 16. GRACIAS
s = slide()
bg2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg2.fill.solid(); bg2.fill.fore_color.rgb = AZUL; bg2.line.fill.background(); bg2.shadow.inherit = False
s.shapes.add_picture(LOGO_COR, Inches(9.9), Inches(4.7), width=Inches(3.6))
text(s, Inches(0.9), Inches(2.7), Inches(10.0), Inches(1.2),
     [{"seg": [("¡GRACIAS!", True)]}], size=44, color=CREMA)
text(s, Inches(0.95), Inches(4.0), Inches(9.5), Inches(0.8),
     [{"seg": [("Quedo a disposición del tribunal para sus preguntas", False)], "italic": True}], size=18, color=RGBColor(0xE9,0xDD,0xCF))
notes(s, "Sonrie, respira, no recojas papeles todavia. Espera las preguntas con calma. "
         "Si no sabes algo, reconduce con honestidad hacia lo que si dominas.")

prs.save(OUT)
print("OK ->", OUT, "| diapositivas:", len(prs.slides._sldIdLst))
